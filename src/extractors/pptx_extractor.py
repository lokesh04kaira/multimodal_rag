from pathlib import Path
from src.utils import clean_text

def extract_pptx(path: str) -> str:
    p = Path(path)
    if not p.exists() or not p.is_file():
        return ""
    try:
        from pptx import Presentation
        prs = Presentation(str(p))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [c.text or "" for c in row.cells]
                            parts.append(" | ".join(cells))
                except Exception:
                    continue
        return clean_text("\n".join(parts))
    except Exception:
        return ""
